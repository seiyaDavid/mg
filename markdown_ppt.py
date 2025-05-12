import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import markdown2

# Convert markdown to HTML first (easier to parse)
md_path = "example.md"  # Replace with your file
with open(md_path, 'r', encoding='utf-8') as f:
    md_content = f.read()

html = markdown2.markdown(md_content)

from bs4 import BeautifulSoup
soup = BeautifulSoup(html, "html.parser")

# Create PowerPoint presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[5]  # Title only or blank

def add_text_slide(text, is_header=False):
    slide = prs.slides.add_slide(blank_slide_layout)
    txBox = slide.shapes.add_textbox(Pt(50), Pt(50), Pt(860), Pt(500))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(32 if is_header else 20)
    if is_header:
        font.bold = True
        font.color.rgb = RGBColor(255, 0, 0)  # Red
    else:
        font.color.rgb = RGBColor(0, 0, 0)  # Black

# Parse the HTML (converted from Markdown)
for tag in soup.find_all(['h1', 'h2', 'h3', 'p']):
    if tag.name.startswith('h'):
        add_text_slide(tag.get_text(strip=True), is_header=True)
    elif tag.name == 'p':
        add_text_slide(tag.get_text(strip=True), is_header=False)

# Save the PowerPoint
output_path = "converted_from_markdown.pptx"
prs.save(output_path)
print(f"Saved presentation to: {output_path}")
