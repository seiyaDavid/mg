from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re

def parse_word_docx_ordered(docx_path):
    doc = Document(docx_path)
    slide_map = {}
    current_slide_num = None
    current_slide = {}
    bullets = []

    for para in doc.paragraphs:
        line = para.text.strip()
        if not line:
            continue

        # Match [SLIDE:X] regardless of position
        slide_match = re.match(r"\[SLIDE:(\d+)\]", line, re.IGNORECASE)
        if slide_match:
            # Save previous slide
            if current_slide_num is not None:
                current_slide['bullets'] = bullets
                slide_map[int(current_slide_num)] = current_slide
                current_slide = {}
                bullets = []
            current_slide_num = int(slide_match.group(1))
            continue

        if line.lower().startswith("heading:"):
            current_slide['heading'] = line.split(":", 1)[1].strip()
        elif line.lower().startswith("bullet:") or line.lower().startswith("subbullet:"):
            bullets.append(line.strip())
        elif line.lower().startswith("note:"):
            current_slide['note'] = line.split(":", 1)[1].strip()

    # Save the last slide
    if current_slide_num is not None:
        current_slide['bullets'] = bullets
        slide_map[int(current_slide_num)] = current_slide

    # Return slide data sorted by slide number
    sorted_slide_data = [slide_map[k] for k in sorted(slide_map.keys())]
    return sorted_slide_data

def create_presentation(slide_data, output_file="from_word.pptx"):
    prs = Presentation()

    for idx, data in enumerate(slide_data):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

        # Remove default placeholders
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

        # Add title with red color and Slide number prefix
        title_text = f"[Slide {idx + 1}] {data.get('heading', '')}"
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.0), Inches(1))
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 0, 0)

        # Add bullets and sub-bullets
        bullets = data.get('bullets', [])
        if bullets:
            body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.0), Inches(4.5))
            tf = body_box.text_frame
            tf.word_wrap = True
            tf.clear()

            for i, bullet in enumerate(bullets):
                level = 1 if bullet.lower().startswith("subbullet:") else 0
                clean_text = re.sub(r"^(Bullet|SubBullet):", "", bullet, flags=re.IGNORECASE).strip()

                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = clean_text
                p.level = level
                p.font.size = Pt(20)
                p.font.name = 'Calibri'

        # Add speaker notes if available
        note_text = data.get('note')
        if note_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = note_text.strip()

    prs.save(output_file)

# Usage
docx_file = "your_input.docx"  # replace with your actual file
slides = parse_word_docx_ordered(docx_file)
create_presentation(slides, output_file="generated_from_word.pptx")
