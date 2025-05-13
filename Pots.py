from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re
from collections import defaultdict

def parse_word_with_tags(word_path):
    doc = Document(word_path)
    slide_data = defaultdict(lambda: {'title': '', 'bullets': [], 'notes': []})

    current_slide = None
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Detect slide tag
        slide_tag_match = re.match(r"\[SLIDE:(\d+)\]", text)
        if slide_tag_match:
            current_slide = int(slide_tag_match.group(1))
            continue

        if current_slide is not None:
            if text.startswith("Heading:"):
                slide_data[current_slide]['title'] = text.replace("Heading:", "").strip()
            elif text.startswith("Bullet:"):
                slide_data[current_slide]['bullets'].append(text.replace("Bullet:", "").strip())
            elif text.startswith("Note:"):
                slide_data[current_slide]['notes'].append(text.replace("Note:", "").strip())

    return slide_data

def create_ppt_from_slide_data(slide_data, output_path):
    prs = Presentation()
    
    for slide_num in sorted(slide_data.keys()):
        data = slide_data[slide_num]
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title text box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.clear()  # Remove default "Click to add..."
        p = title_frame.paragraphs[0]
        run = p.add_run()
        run.text = data['title']
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 0, 0)  # Red

        # Add bullet points text box
        if data['bullets']:
            bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(4))
            bullet_frame = bullet_box.text_frame
            bullet_frame.word_wrap = True
            for bullet in data['bullets']:
                p = bullet_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = Pt(6)
                p.font.size = Pt(20)

        # Add notes
        if data['notes']:
            notes = slide.notes_slide.notes_text_frame
            notes.clear()
            notes.text = '\n'.join(data['notes'])

    prs.save(output_path)
    print(f"✅ PowerPoint saved to: {output_path}")

# Example usage
slide_data = parse_word_with_tags("your_doc.docx")
create_ppt_from_slide_data(slide_data, "output_presentation.pptx")
